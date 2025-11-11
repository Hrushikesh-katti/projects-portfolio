import pdfplumber
from docx import Document
from pptx import Presentation
import spacy
import nltk
from nltk.corpus import stopwords
import string
import re
import os
from fuzzywuzzy import fuzz
from openai import OpenAI
import streamlit as st
import time
import tempfile
from streamlit.components.v1 import html

# Download NLTK stopwords and punkt if not already downloaded
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')
try:
    nltk.data.find('tokenizers/punkt_tab')
except LookupError:
    nltk.download('punkt_tab')

# Load SpaCy model
nlp = spacy.load('en_core_web_sm')

# Predefined list of common skills
SKILL_LIST = [
    'python', 'sql', 'excel', 'power bi', 'tableau', 'django', 'html', 'css', 'javascript', 'react',
    'adobe photoshop', 'illustrator', 'figma', 'motion graphics', 'grammarly', 'semrush', 'surferseo',
    'wordpress', 'canva', 'agile', 'scrum', 'budgeting', 'leadership', 'jira', 'team management',
    'risk assessment', 'communication', 'seo', 'sem', 'google ads', 'meta ads', 'content creation',
    'public speaking', 'social media', 'penetration testing', 'firewall', 'iso 27001', 'gdpr',
    'recruitment', 'scheduling', 'crm', 'field sales', 'tensorflow', 'nlp', 'vision', 'data analysis',
    'project management', 'web development', 'microservices', 'rest apis', 'problem solving', 'analytics',
    'agile methodology'
]

# OpenRouter API configuration for DeepSeek
OPENROUTER_API_KEY = "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"  # Replace with your actual API key
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY
)

# ================== BACKEND FUNCTIONS ==================
def parse_resume(content):
    try:
        if isinstance(content, str) and content.startswith('http'):
            return "URL parsing not supported yet. Please use text or file upload."
        elif isinstance(content, str):
            return content.strip()
        elif content.name.endswith('.pdf'):
            with pdfplumber.open(content) as pdf:
                text = ''
                for page in pdf.pages:
                    text += page.extract_text() or ''
                return text.strip()
        elif content.name.endswith('.txt'):
            return content.read().decode('utf-8').strip()
        elif content.name.endswith('.docx'):
            doc = Document(content)
            return '\n'.join(para.text for para in doc.paragraphs).strip()
        elif content.name.endswith('.pptx'):
            prs = Presentation(content)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
            return text.strip()
        else:
            return "Unsupported file format. Please use PDF, .txt, .docx, or .pptx."
    except Exception as e:
        return f"Error parsing file: {str(e)}"

def preprocess_text(text):
    text = re.sub(r'CONFIDENTIAL:.*?clauses\.', '', text, flags=re.DOTALL)
    text = re.sub(r'CCOONNFFIIDDEENNTTIIAALL.*?ccllaauusseess\.\.', '', text, flags=re.DOTALL)
    doc = nlp(text.lower())
    tokens = [token for token in doc if token.ent_type_ != 'PERSON' and token.text not in ['name', 'role', 'resume', 'experience', 'year', 'years']]
    stop_words = set(stopwords.words('english'))
    tokens = [token.lemma_ for token in tokens if token.text not in string.punctuation and token.text not in stop_words]
    return ' '.join(tokens)

def extract_experience(text):
    patterns = [r'(\d+)\s*(?:year|years)\s*(?:of\s*experience)?', r'(\d+)\s*\+\s*(?:year|years)\s*(?:of\s*experience)?', r'fresher']
    for pattern in patterns:
        match = re.search(pattern, text.lower())
        if match:
            if match.group(0) == 'fresher':
                return 0
            return int(match.group(1))
    return 0

def check_eligibility(resume_text, jd_text):
    resume_experience = extract_experience(resume_text)
    jd_experience = extract_experience(jd_text)
    return "Eligible" if resume_experience >= jd_experience else "Not Eligible"

def extract_keywords(text):
    processed_text = preprocess_text(text)
    doc = nlp(processed_text)
    keywords = [chunk.text.lower() for chunk in doc.noun_chunks if len(chunk.text.split()) <= 2] + [token.text.lower() for token in doc if token.pos_ in ['NOUN', 'PROPN'] and len(token.text) > 2]
    filtered_keywords = [skill for kw in set(keywords) for skill in SKILL_LIST if fuzz.partial_ratio(kw, skill) > 75 or skill in kw.split()]
    return list(set(filtered_keywords))

def optimize_resume(resume_text, jd_text, matching_skills, missing_skills):
    prompt = f"Optimize the following resume to align with the job description. Emphasize matching skills: {matching_skills}. For missing skills: {missing_skills}, add sections or rephrase to suggest familiarity or willingness to learn. Keep it professional, concise, under 500 words. Return only the optimized resume.\n\nResume:\n{resume_text}\n\nJob Description:\n{jd_text}"
    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-r1-0528-qwen3-8b:free",
            messages=[{"role": "system", "content": "You are a professional resume writer."}, {"role": "user", "content": prompt}],
            temperature=0.7, top_p=0.9, max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error optimizing resume: {str(e)}"

def generate_interview_guidance(jd_text, matching_skills, missing_skills):
    prompt = f"Generate interview guidance for the job description. Leverage matching skills: {matching_skills}. Address missing skills: {missing_skills} with strategies (e.g., willingness to learn). Include 3-5 relevant questions with brief answers. Keep it concise, under 300 words.\n\nJob Description:\n{jd_text}"
    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-r1-0528-qwen3-8b:free",
            messages=[{"role": "system", "content": "You are a career coach."}, {"role": "user", "content": prompt}],
            temperature=0.7, top_p=0.9, max_tokens=500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error generating guidance: {str(e)}"

def compute_match_score(resume_text, jd_text):
    eligibility = check_eligibility(resume_text, jd_text)
    resume_keywords = extract_keywords(resume_text)
    jd_keywords = extract_keywords(jd_text)
    matching_skills = [s for s in resume_keywords for j in jd_keywords if fuzz.partial_ratio(s, j) > 75 or s == j]
    matching_skills = list(set(matching_skills))
    skillset_match_score = 100 * (len(matching_skills) / max(len(jd_keywords), 1))
    skillset_match_score = round(skillset_match_score, 2)
    if not matching_skills:
        skillset_match_score = 0.0
    if eligibility == "Not Eligible":
        skillset_match_score = min(skillset_match_score, 5.0)
    missing_skills = [skill for skill in jd_keywords if skill not in matching_skills]
    extra_skills = [skill for skill in resume_keywords if skill not in matching_skills]
    return {'eligibility': eligibility, 'skillset_match_score': skillset_match_score, 'matching_skills': matching_skills, 'missing_skills': missing_skills, 'extra_skills': extra_skills}

def process_resume_and_jd(resume_input, jd_input):
    with st.spinner("Crunching the data... summoning the resume gods 🔮"):
        time.sleep(1)  # Simulate processing time
        if not resume_input or not jd_input:
            return None, "Please provide both resume and job description.", "", ""
        
        if hasattr(resume_input, 'read'):
            resume_text = parse_resume(resume_input)
        else:
            resume_text = parse_resume(resume_input)
            
        if resume_text.startswith("Error") or resume_text.startswith("Unsupported"):
            return None, resume_text, "", ""
            
        if hasattr(jd_input, 'read'):
            jd_text = parse_resume(jd_input)
        else:
            jd_text = jd_input
            
        if jd_text.startswith("Error") or jd_text.startswith("Unsupported"):
            return None, f"Error reading JD: {jd_text}", "", ""
            
        result = compute_match_score(resume_text, jd_text)
        match_results = (f"Eligibility: {result['eligibility']}\n"
                         f"Skillset Match Score: {result['skillset_match_score']}%\n"
                         f"Matching Skills: {result['matching_skills']}\n"
                         f"Missing Skills: {result['missing_skills']}\n"
                         f"Extra Skills: {result['extra_skills']}")
        optimized_resume = ""
        interview_guidance = ""
        if result['eligibility'] == "Eligible" and result['skillset_match_score'] > 0:
            optimized_resume = optimize_resume(resume_text, jd_text, result['matching_skills'], result['missing_skills'])
            interview_guidance = generate_interview_guidance(jd_text, result['matching_skills'], result['missing_skills'])
        return result, match_results, optimized_resume, interview_guidance

# ================== STREAMLIT UI ==================
st.set_page_config(page_title="JobFit AI", layout="wide", page_icon="💼")

# Dark Theme CSS Styling with Orange Accents
st.markdown("""
<style>
    /* Main app background */
    .stApp {
        background: #121212 !important;
        color: #ffffff !important;
    }
    
    /* All text color */
    body, .stTextInput>div>div>input, .stTextArea>div>div>textarea, 
    .stSelectbox>div>div>select, .stMarkdown, .stAlert, 
    .stButton>button, .stFileUploader>div>div>div>p {
        color: #ffa500 !important;
    }
    
    /* Card styling */
    .card {
        border-radius: 12px !important;
        box-shadow: 0 6px 18px rgba(0, 0, 0, 0.3) !important;
        padding: 25px !important;
        background: #1e1e1e !important;
        margin-bottom: 25px !important;
        border: 1px solid #333 !important;
    }
    
    /* Modern button */
    .stButton>button {
        background: linear-gradient(45deg, #ff8c00, #ff4500) !important;
        border: none !important;
        color: #121212 !important;
        padding: 12px 24px !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 16px !important;
        transition: all 0.3s !important;
        box-shadow: 0 4px 6px rgba(255, 140, 0, 0.3) !important;
    }
    
    .stButton>button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 6px 12px rgba(255, 140, 0, 0.4) !important;
    }
    
    /* Progress spinner */
    .stSpinner>div {
        background: linear-gradient(45deg, #ff8c00, #ff4500) !important;
    }
    
    /* Custom header */
    .header {
        background: linear-gradient(90deg, #1a1a1a 0%, #333 100%) !important;
        padding: 2rem !important;
        color: #ffa500 !important;
        border-radius: 0 0 12px 12px !important;
        margin-bottom: 2rem !important;
        text-align: center !important;
        border-bottom: 3px solid #ff8c00 !important;
    }
    
    /* Skill tags */
    .skill-tag {
        display: inline-block !important;
        padding: 4px 12px !important;
        background: #333 !important;
        border-radius: 16px !important;
        margin: 4px !important;
        font-size: 14px !important;
        color: #ffa500 !important;
        border: 1px solid #ff8c00 !important;
    }
    
    .missing-skill-tag {
        display: inline-block !important;
        padding: 4px 12px !important;
        background: #330000 !important;
        border-radius: 16px !important;
        margin: 4px !important;
        font-size: 14px !important;
        color: #ff6666 !important;
        border: 1px solid #ff3333 !important;
    }
    
    .extra-skill-tag {
        display: inline-block !important;
        padding: 4px 12px !important;
        background: #003300 !important;
        border-radius: 16px !important;
        margin: 4px !important;
        font-size: 14px !important;
        color: #66ff66 !important;
        border: 1px solid #00cc00 !important;
    }
    
    .match-percentage {
        font-size: 3rem !important;
        font-weight: 800 !important;
        background: linear-gradient(45deg, #ff8c00, #ff4500) !important;
        -webkit-background-clip: text !important;
        -webkit-text-fill-color: transparent !important;
        text-align: center !important;
        margin: 10px 0 !important;
    }
    
    /* Text areas */
    .stTextArea>div>div>textarea {
        border-radius: 8px !important;
        padding: 12px !important;
        background: #1e1e1e !important;
        color: #ffa500 !important;
        border: 1px solid #ff8c00 !important;
    }
    
    /* File uploader */
    .stFileUploader>div>div {
        border-radius: 8px !important;
        padding: 12px !important;
        border: 1px dashed #ff8c00 !important;
        background: #1e1e1e !important;
    }
    
    /* Tabs */
    .stTabs>div>div>button {
        border-radius: 8px 8px 0 0 !important;
        padding: 8px 16px !important;
        background: #1e1e1e !important;
        color: #ffa500 !important;
        border: 1px solid #333 !important;
    }
    
    .stTabs>div>div>button[aria-selected="true"] {
        background: #333 !important;
        font-weight: bold !important;
        color: #ff8c00 !important;
        border-bottom: 2px solid #ff8c00 !important;
    }
    
    /* Container padding */
    .stContainer {
        padding: 1rem !important;
    }
    
    /* Custom hr */
    .custom-hr {
        border: 0 !important;
        height: 1px !important;
        background: linear-gradient(90deg, transparent, #ff8c00, transparent) !important;
        margin: 2rem 0 !important;
    }
    
    /* Fix for markdown text color */
    .stMarkdown p, .stMarkdown li, .stMarkdown ol, .stMarkdown ul {
        color: #ffa500 !important;
    }
    
    /* Placeholder text color */
    .stTextArea>div>div>textarea::placeholder {
        color: #666 !important;
    }
    
    /* Input labels */
    .stTextInput label, .stTextArea label, .stFileUploader label {
        color: #ffa500 !important;
    }
    
    /* Caption text */
    .stCaption {
        color: #888 !important;
    }
</style>
""", unsafe_allow_html=True)

# Header Section
st.markdown("""
<div class='header'>
    <h1 style="margin: 0; color: #ffa500;">JobFit AI</h1>
    <p style="margin: 0; opacity: 0.9; font-size: 1.2rem; color: #ffa500;">Resume Optimization & Interview Preparation Tool</p>
    <p style="margin: 0; font-size: 0.8rem; opacity: 0.7; color: #ffa500;">AI-powered career assistant</p>
</div>
""", unsafe_allow_html=True)

# Main Content
with st.container():
    col1, col2 = st.columns(2, gap="large")

    with col1:
        st.markdown("### 📄 Resume Input")
        with st.container():
            resume_text = st.text_area(
                "Paste your resume text here", 
                height=200,
                placeholder="Or upload a file below...",
                key="resume_text"
            )
            resume_file = st.file_uploader(
                "Upload Resume (PDF, DOCX, PPTX, TXT)",
                type=["pdf", "docx", "pptx", "txt"],
                key="resume_file"
            )
            st.caption("Supported formats: PDF, DOCX, PPTX, TXT")

    with col2:
        st.markdown("### 📋 Job Description Input")
        with st.container():
            jd_text = st.text_area(
                "Paste the job description here", 
                height=200,
                placeholder="Or upload a text file...",
                key="jd_text"
            )
            jd_file = st.file_uploader(
                "Upload Job Description (TXT)",
                type=["txt"],
                key="jd_file"
            )
            st.caption("Supported format: TXT")

# Determine input sources
resume_input = resume_file if resume_file else resume_text
jd_input = jd_file if jd_file else jd_text

# Process Button
col1, col2, col3 = st.columns([1, 2, 1])
with col2:
    analyze_btn = st.button(
        "🚀 Analyze & Optimize", 
        use_container_width=True,
        type="primary"
    )

# Results Section
if analyze_btn:
    result, match_results, optimized_resume, interview_guidance = process_resume_and_jd(resume_input, jd_input)
    
    if result is None:
        st.error(match_results)
    else:
        st.markdown('<div class="custom-hr"></div>', unsafe_allow_html=True)
        
        # Results in Tabs
        tab1, tab2, tab3 = st.tabs(["📊 Match Analysis", "📝 Optimized Resume", "🎤 Interview Prep"])
        
        with tab1:
            col1, col2 = st.columns([1, 2])

            with col1:
                st.markdown(f"""
                <div class="card" style="text-align: center;">
                    <div class="match-percentage">{result['skillset_match_score']}%</div>
                    <div style="font-size: 1.2rem; color: #ff8c00; text-align: center;">Match Score</div>
                    <div style="margin-top: 10px; font-size: 0.9rem; text-align: center;">
                        Status: <span style="color: {'#4CAF50' if result['eligibility'] == 'Eligible' else '#F44336'}">
                        {result['eligibility']}</span>
                    </div>
                </div>
                """, unsafe_allow_html=True)

            with col2:
                st.markdown("""
                <div class="card">
                    <h3 style="margin-top: 0; color: #ffa500;">Skills Analysis</h3>
                """, unsafe_allow_html=True)

                # Matching Skills
                matching_skills_html = ' '.join([f'<span class="skill-tag">{skill}</span>' for skill in result['matching_skills']])
                st.markdown(f"""
                    <h4 style="color: #4CAF50;">✅ Matching Skills ({len(result['matching_skills'])})</h4>
                    <div>{matching_skills_html}</div>
                """, unsafe_allow_html=True)

                # Missing Skills
                missing_skills_html = ' '.join([f'<span class="missing-skill-tag">{skill}</span>' for skill in result['missing_skills']])
                st.markdown(f"""
                    <h4 style="margin-top: 15px; color: #F44336;">⚠️ Missing Skills ({len(result['missing_skills'])})</h4>
                    <div>{missing_skills_html}</div>
                """, unsafe_allow_html=True)

                # Extra Skills
                extra_skills_html = ' '.join([f'<span class="extra-skill-tag">{skill}</span>' for skill in result['extra_skills']])
                st.markdown(f"""
                    <h4 style="margin-top: 15px; color: #66ff66;">➕ Extra Skills ({len(result['extra_skills'])})</h4>
                    <div>{extra_skills_html}</div>
                </div>
                """, unsafe_allow_html=True)

        
        with tab2:
            if optimized_resume:
                st.markdown(f"""
                <div class="card">
                    <h3 style="margin-top: 0; color: #ffa500;">Optimized Resume</h3>
                    <div style="white-space: pre-wrap; line-height: 1.6; color: #ffa500;">{optimized_resume}</div>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.warning("No optimized resume generated due to low match score or ineligibility")
        
        with tab3:
            if interview_guidance:
                st.markdown(f"""
                <div class="card">
                    <h3 style="margin-top: 0; color: #ffa500;">Interview Preparation Guide</h3>
                    <div style="white-space: pre-wrap; line-height: 1.6; color: #ffa500;">{interview_guidance}</div>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.warning("No interview guidance generated due to low match score or ineligibility")

# Footer
st.markdown("""
<div style="text-align: center; margin-top: 50px; color: #888; font-size: 0.8rem;">
    <p>JobFit AI - Powered by NLP and DeepSeek</p>
    <p>Note: AI-generated content should be reviewed before use</p>
</div>
""", unsafe_allow_html=True)